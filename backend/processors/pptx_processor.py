"""
PPTX document processor
Extracts text from PowerPoint presentations
"""
from typing import Dict, Any
from pathlib import Path
from pptx import Presentation
from backend.processors.text_chunker import get_text_chunker


class PPTXProcessor:
    """Processes PPTX presentations"""
    
    def __init__(self):
        self.chunker = get_text_chunker()
    
    def process(self, file_path: Path) -> Dict[str, Any]:
        """
        Process PPTX file
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            Dictionary with chunks and metadata
        """
        try:
            prs = Presentation(file_path)
            
            all_chunks = []
            chunk_metadatas = []
            total_slides = len(prs.slides)
            
            # Extract text from each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text_parts = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text)
                
                # Combine slide text
                if slide_text_parts:
                    slide_text = "\n".join(slide_text_parts)
                    
                    # Chunk the slide text
                    chunks = self.chunker.chunk_text(slide_text)
                    
                    # Create metadata for each chunk
                    for chunk in chunks:
                        all_chunks.append(chunk)
                        chunk_metadatas.append({
                            "filename": file_path.name,
                            "file_type": "pptx",
                            "page": slide_num,  # Using slide number as "page"
                            "total_pages": total_slides
                        })
            
            return {
                "chunks": all_chunks,
                "metadatas": chunk_metadatas,
                "total_pages": total_slides,
                "total_chunks": len(all_chunks)
            }
            
        except Exception as e:
            raise RuntimeError(f"PPTX processing failed: {str(e)}")


def get_pptx_processor() -> PPTXProcessor:
    """Get PPTXProcessor instance"""
    return PPTXProcessor()
